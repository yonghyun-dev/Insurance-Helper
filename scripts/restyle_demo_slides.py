"""
데모 2장(의미검색 / AI 챗봇)을 제안서 12·13페이지 스타일로 교체한다.

파일 경로: scripts/restyle_demo_slides.py
목적:
    현재 파일(슬라이드 1 아키텍처 + 2 의미검색 + 3 챗봇)에서 2·3페이지를 삭제하고,
    "상단 제목 + 2줄 부연설명 / 좌측 실제 화면 / 우측 구현 방안(자세하지만 쉬운 설명)" 구조로
    다시 만든다. 슬라이드 1은 그대로 둔다.
주의:
    현재 파일을 직접 수정한다(.bak5 백업 별도). 백업에서 복원하지 않는다(사용자가 지운 페이지 보존).
주요 의존성:
    python-pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = r"C:\Users\edgar\Desktop\데이터카탈로그 소개서\VODA_AI_Architecture.pptx"
CACHE = r"C:\Users\edgar\.claude\image-cache\c1de5183-685f-405c-bc75-472b2970def3"
SHOT7 = os.path.join(CACHE, "7.png")  # 전체 카탈로그
SHOT10 = os.path.join(CACHE, "10.png")  # 의미검색 결과
SHOT11 = os.path.join(CACHE, "11.png")  # AI 챗봇 대화
IMG_AR = 866 / 1815  # 카탈로그/결과 가로형 비율(h/w)
CHAT_AR = 1448 / 1086  # 챗봇 세로형 비율
IMG_W = 1815

NAVY = RGBColor(0x1F, 0x2D, 0x50)
BLUE = RGBColor(0x2B, 0x5F, 0xC0)
GRAY = RGBColor(0x55, 0x5E, 0x6E)
LIGHTGRAY = RGBColor(0x9A, 0xA2, 0xB0)
RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BD = RGBColor(0xB8, 0xC4, 0xDB)
FONT = "맑은 고딕"


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def box(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    anchor=MSO_ANCHOR.MIDDLE,
    align=PP_ALIGN.LEFT,
    round_=False,
    wrap=True,
    space_after=2,
):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    b = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    b.shadow.inherit = False
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(line_w)
    tf = b.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        text, size, color, bold = ln
        r = p.add_run()
        set_run(r, text, size, color, bold)
    return b


def oval_ring(slide, cx, cy, d, color):
    o = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d)
    )
    o.shadow.inherit = False
    o.fill.background()
    o.line.color.rgb = color
    o.line.width = Pt(2.5)
    return o


def seg(slide, x1, y1, x2, y2, color, width=2.25):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def picture(slide, path, x, y, w, border):
    if os.path.exists(path):
        p = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
        p.line.color.rgb = border
        p.line.width = Pt(1.25)
        return p
    box(
        slide,
        x,
        y,
        w,
        w * 0.6,
        [("(이미지 없음)", 11, LIGHTGRAY, False)],
        fill=PANEL_BD,
        line=border,
    )
    return None


def header(slide, title, subtitle):
    box(
        slide,
        0.5,
        0.18,
        12.3,
        0.52,
        [(title, 20, NAVY, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    box(slide, 0.52, 0.72, 3.2, 0.055, [], fill=BLUE)
    box(
        slide,
        0.5,
        0.82,
        12.4,
        0.6,
        [(subtitle, 12.5, GRAY, False)],
        anchor=MSO_ANCHOR.TOP,
        align=PP_ALIGN.LEFT,
        wrap=True,
    )


def impl_panel(slide, blocks):
    px, pw = 8.15, 4.7
    box(
        slide,
        px,
        1.55,
        pw,
        0.48,
        [("구현 방안", 14, WHITE, True)],
        fill=BLUE,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        round_=True,
    )
    ys = [2.3, 3.85, 5.4]
    for n, ((title, subs), y) in enumerate(zip(blocks, ys), 1):
        box(
            slide,
            px,
            y,
            0.36,
            0.36,
            [(str(n), 13, WHITE, True)],
            fill=BLUE,
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
            round_=True,
        )
        lines = [(title, 12.5, NAVY, True)] + [("• " + sb, 10.5, GRAY, False) for sb in subs]
        box(
            slide,
            px + 0.46,
            y - 0.05,
            pw - 0.5,
            1.45,
            lines,
            anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT,
            space_after=3,
        )


def footnote(slide, w=7.5):
    box(
        slide,
        0.4,
        6.6,
        w,
        0.3,
        [("※ 화면은 예시이며 실제 UI와 다를 수 있습니다.", 9.5, LIGHTGRAY, False)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def find_blank_layout(prs):
    best, best_n = prs.slide_layouts[0], len(prs.slide_layouts[0].placeholders)
    for lay in prs.slide_layouts:
        if len(lay.placeholders) < best_n:
            best, best_n = lay, len(lay.placeholders)
    return best


def clear_shapes(slide):
    """슬라이드의 모든 도형을 제거(슬라이드 자체는 유지) — 내용만 새로 그린다."""
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def build_semantic(s):
    header(
        s,
        "의미(Semantic) 검색 — 말의 '뜻'으로 데이터 찾기",
        "정확한 데이터 이름이나 키워드를 몰라도, 평소 쓰는 말로 검색하면 그 의미에 맞는 데이터를 "
        "AI가 찾아주어 검색 정확도를 높입니다.",
    )
    # 좌: 의미검색 결과 캡처
    box(
        s,
        0.4,
        1.55,
        7.4,
        0.32,
        [("「배로 실어 나른 물량」 검색 결과 — GLOVE 물류 실적 4건", 12, NAVY, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    picture(s, SHOT10, 0.4, 1.95, 7.4, BLUE)
    # 우: 구현 방안
    impl_panel(
        s,
        [
            (
                "키워드 + 의미(시멘틱) 하이브리드 검색",
                [
                    "정확한 데이터 이름을 몰라도, 입력한 말의 '뜻'으로 찾아줍니다",
                    "단어가 똑같지 않아도(동의어·비슷한 말) 관련 데이터를 함께 보여줍니다",
                ],
            ),
            (
                "정확도를 높이는 똑똑한 정렬",
                [
                    "찾은 결과를 질문과 가까운 순서로 다시 정렬해 위에서부터 보여줍니다",
                    "현업에서 쓰는 용어도 알아듣도록 '비슷한 말 사전'을 함께 활용합니다",
                ],
            ),
            (
                "권한에 맞는 안전한 검색 + 바로 활용",
                [
                    "내 권한(부문·보안등급)에 맞는 데이터만 검색·노출됩니다",
                    "찾은 데이터를 선택하면 곧바로 리포트 작성 화면으로 이어집니다",
                ],
            ),
        ],
    )
    footnote(s, 7.5)
    return s


def build_chatbot(s):
    header(
        s,
        "AI 챗봇 — 물어보면 찾아주는 데이터 도우미",
        "모든 화면에서 플로팅 버튼으로 챗봇을 열어 평소 말투로 물어보면, 적절한 데이터를 찾아 답하고 "
        "더 나은 질문을 하도록 가이드합니다.",
    )
    # 좌: 카탈로그(플로팅 버튼) + 챗봇 대화 캡처
    ix, iy, iw = 0.4, 1.95, 4.0
    sc = iw / IMG_W
    picture(s, SHOT7, ix, iy, iw, PANEL_BD)
    bx, by = ix + 1740 * sc, iy + 795 * sc
    oval_ring(s, bx, by, 0.46, RED)
    box(
        s,
        0.5,
        iy + 2.0,
        4.0,
        0.3,
        [("↑ 모든 화면 우하단 '플로팅 챗봇 버튼'", 10, RED, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    cx0, cy0, cw = 4.75, 1.65, 2.95
    picture(s, SHOT11, cx0, cy0, cw, PANEL_BD)
    seg(s, bx + 0.25, by - 0.05, cx0 - 0.05, cy0 + 1.0, RED)
    # 우: 구현 방안
    impl_panel(
        s,
        [
            (
                "어디서나 쉽게 — 플로팅 버튼",
                [
                    "모든 화면 우측 하단 챗봇 버튼으로, 보던 화면을 떠나지 않고 질문",
                    "답을 받은 뒤 '검색결과로 이동'으로 곧장 결과 화면까지 (끊김 없는 경험)",
                ],
            ),
            (
                "대화로 좁혀가는 똑똑한 안내",
                [
                    "질문이 모호하면 챗봇이 되물어 원하는 것을 구체화하도록 안내합니다",
                    "업무용어(예: '글로벌 물류')를 실제 시스템명(GLOVE)으로 알아서 매핑",
                ],
            ),
            (
                "점점 똑똑해지는 챗봇 (확장)",
                ["메뉴별 사용법 안내, 자주 묻는 질문 대응 등 활용 범위 확대를 검토합니다"],
            ),
        ],
    )
    footnote(s, 7.5)
    return s


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)
    # 슬라이드 2(의미검색)·3(챗봇)의 도형만 비우고 새로 그린다 — 슬라이드 1은 유지
    clear_shapes(slides[1])
    build_semantic(slides[1])
    clear_shapes(slides[2])
    build_chatbot(slides[2])
    prs.save(SRC)
    print("saved:", SRC)
    print("total slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
