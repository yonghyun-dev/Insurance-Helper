"""
VODA 데이터 카탈로그 검색 — 유스 케이스 & 시나리오 개체관계도 1장을 PPT에 추가한다.

파일 경로: scripts/add_voda_scenarios.py
목적:
    기존 아키텍처 슬라이드(1장) 뒤에, 두 검색 시나리오를 하나의 개체관계도로 담은 장표를 추가한다.
    - 사용자(Actor)는 왼쪽에 따로 분리
    - ① 카탈로그 검색 / ② 챗봇 검색 / 공통 RAG·거버넌스 세 영역으로 구분
    - 박스는 제목 한 줄만 (작은 설명글 없음)
주요 의존성:
    python-pptx
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = r"C:\Users\edgar\Desktop\데이터카탈로그 소개서\VODA_AI_Architecture.pptx"

NAVY = RGBColor(0x2A, 0x36, 0x56)
BLUE = RGBColor(0x2B, 0x5F, 0xC0)  # 시나리오 ① 경로
ORANGE = RGBColor(0xE0, 0x6A, 0x18)  # 시나리오 ② 경로
GRAYLINE = RGBColor(0x6E, 0x78, 0x8C)  # 공통 경로
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5A, 0x64, 0x78)

# 영역(밴드) 색
BL_FILL = RGBColor(0xEA, 0xF1, 0xFB)
BL_BD = RGBColor(0xA9, 0xC2, 0xE8)
OR_FILL = RGBColor(0xFB, 0xEF, 0xE3)
OR_BD = RGBColor(0xE0, 0xA8, 0x78)
GY_FILL = RGBColor(0xF1, 0xF3, 0xF6)
GY_BD = RGBColor(0xC4, 0xCB, 0xD6)
FONT = "맑은 고딕"

NODES = {}


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_text(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    anchor=MSO_ANCHOR.MIDDLE,
    align=PP_ALIGN.CENTER,
    round_=False,
    wrap=True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.shadow.inherit = False
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        text, size, color, bold = ln
        r = p.add_run()
        set_run(r, text, size, color, bold)
    return box


def reg(key, cx, cy, w, h):
    NODES[key] = (cx - w / 2, cy - h / 2, w, h)


def draw_node(slide, key, label):
    x, y, w, h = NODES[key]
    add_text(
        slide,
        x,
        y,
        w,
        h,
        [(label, 12, NAVY, True)],
        fill=WHITE,
        line=NAVY,
        line_w=1.5,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
    )


def right(k):
    x, y, w, h = NODES[k]
    return (x + w, y + h / 2)


def left(k):
    x, y, w, h = NODES[k]
    return (x, y + h / 2)


def top(k):
    x, y, w, h = NODES[k]
    return (x + w / 2, y)


def bottom(k):
    x, y, w, h = NODES[k]
    return (x + w / 2, y + h)


def _seg(slide, x1, y1, x2, y2, color, width, end_arrow=False, begin_arrow=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if begin_arrow:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if end_arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def connect(slide, pts, color, label="", lx=0, ly=0, lw=1.0, *, begin_arrow=False, width=1.75):
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        _seg(
            slide,
            x1,
            y1,
            x2,
            y2,
            color,
            width,
            end_arrow=(i == len(pts) - 2),
            begin_arrow=(begin_arrow and i == 0),
        )
    if label:
        add_text(
            slide,
            lx - lw / 2,
            ly - 0.13,
            lw,
            0.26,
            [(label, 9, NAVY, False)],
            fill=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
            wrap=False,
        )


def band(slide, x, y, w, h, label, fill, border, label_color):
    add_text(slide, x, y, w, h, [], fill=fill, line=border, line_w=1.5, round_=True)
    add_text(
        slide,
        x + 0.12,
        y + 0.1,
        3.2,
        0.34,
        [(label, 12.5, label_color, True)],
        anchor=MSO_ANCHOR.TOP,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def find_blank_layout(prs):
    best, best_n = prs.slide_layouts[0], len(prs.slide_layouts[0].placeholders)
    for lay in prs.slide_layouts:
        if len(lay.placeholders) < best_n:
            best, best_n = lay, len(lay.placeholders)
    return best


def clear_placeholders(slide):
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def build_diagram(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)

    # 제목 밴드
    add_text(s, 0, 0, 13.333, 0.9, [], fill=NAVY)
    add_text(
        s,
        0.35,
        0.05,
        12.6,
        0.8,
        [("데이터 카탈로그 검색 — 유스 케이스 & 시나리오 개체관계도", 20, WHITE, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )

    # 개체 좌표 등록
    reg("user", 1.2, 4.0, 2.2, 0.72)
    reg("cat", 3.95, 2.5, 2.3, 0.72)
    reg("result", 3.95, 3.7, 2.3, 0.72)
    reg("chat", 3.95, 5.6, 2.3, 0.72)
    reg("agent", 7.5, 2.9, 2.4, 0.72)
    reg("rag", 7.5, 4.3, 2.4, 0.72)
    reg("vdb", 7.5, 5.7, 2.6, 0.72)
    reg("extract", 10.9, 4.3, 2.5, 0.72)
    reg("catalog", 10.9, 5.7, 2.5, 0.72)

    # 영역(밴드) — 박스/선보다 먼저 그려 뒤에 깔리게
    band(s, 2.65, 1.55, 2.85, 2.75, "① 카탈로그 검색", BL_FILL, BL_BD, BLUE)
    band(s, 2.65, 4.95, 2.85, 1.35, "② 챗봇 검색", OR_FILL, OR_BD, ORANGE)
    band(s, 6.35, 1.55, 6.1, 4.85, "공통 — AI 서비스(RAG) · 거버넌스", GY_FILL, GY_BD, GRAY)

    # 관계선
    # 사용자 → 두 진입점 (색만으로 구분, 라벨 없음)
    connect(s, [right("user"), (2.5, 4.0), (2.5, 2.5), left("cat")], BLUE)
    connect(s, [right("user"), (2.5, 4.0), (2.5, 5.6), left("chat")], ORANGE)

    # 시나리오 ① (파랑)
    connect(
        s,
        [right("cat"), (5.75, 2.5), (5.75, 2.9), left("agent")],
        BLUE,
        "1..* 질의",
        5.75,
        2.68,
        0.85,
    )
    connect(
        s,
        [(left("agent")[0], 3.0), (5.95, 3.0), (5.95, 3.7), right("result")],
        BLUE,
        "① 결과 반환",
        5.95,
        3.35,
        1.05,
    )

    # 시나리오 ② (주황)
    connect(
        s,
        [right("chat"), (6.2, 5.6), (6.2, 2.9), left("agent")],
        ORANGE,
        "질의 / 답변",
        6.2,
        5.0,
        0.95,
        begin_arrow=True,
    )

    # 공통 RAG · 거버넌스 (회색)
    connect(s, [bottom("agent"), top("rag")], GRAYLINE, "1..* RAG 검색", 7.5, 3.6, 1.35)
    connect(s, [bottom("rag"), top("vdb")], GRAYLINE, "Knowledge 조회", 7.5, 5.0, 1.4)
    connect(s, [bottom("extract"), top("catalog")], GRAYLINE, "메타 적재", 10.9, 5.0, 1.1)
    connect(
        s,
        [left("catalog"), right("vdb")],
        GRAYLINE,
        "인덱싱",
        (right("vdb")[0] + left("catalog")[0]) / 2,
        5.7,
        0.8,
    )

    # 개체 박스 (선 위에)
    draw_node(s, "user", "사용자 (Creator/현업)")
    draw_node(s, "cat", "카탈로그 검색 화면")
    draw_node(s, "result", "검색결과 페이지")
    draw_node(s, "chat", "AI 서비스 채팅창")
    draw_node(s, "agent", "AI Agent")
    draw_node(s, "rag", "RAG 검색")
    draw_node(s, "vdb", "Knowledge / Vector DB")
    draw_node(s, "extract", "메타 추출 · 인덱싱")
    draw_node(s, "catalog", "메타 카탈로그")
    return s


def main():
    prs = Presentation(SRC)
    layout = find_blank_layout(prs)
    build_diagram(prs, layout)
    prs.save(SRC)
    print("saved:", SRC)
    print("total slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
