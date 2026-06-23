"""
고객 소개용 데모 장표 2장(의미검색 / AI 챗봇)을 기존 PPT에 이어붙인다.

파일 경로: scripts/add_demo_slides.py
목적:
    실제 카탈로그 화면 캡처를 메인으로 깔아 고객이 바로 이해하는 데모 장표를 만든다.
    - 5p 의미검색: 전체 카탈로그(7.png) → 일상어 검색 → 의미검색 결과(10.png) Before/After
    - 6p AI 챗봇: 카탈로그 + 플로팅 버튼 강조 + 챗봇 핑퐁 대화 말풍선 (제안서 12p 톤)
주의:
    .bak4(4슬라이드) 복원 후 실행 → 5·6페이지로 추가.
주요 의존성:
    python-pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = r"C:\Users\edgar\Desktop\데이터카탈로그 소개서\VODA_AI_Architecture.pptx"
CACHE = r"C:\Users\edgar\.claude\image-cache\c1de5183-685f-405c-bc75-472b2970def3"
SHOT7 = os.path.join(CACHE, "7.png")  # 전체 카탈로그
SHOT10 = os.path.join(CACHE, "10.png")  # 의미검색 결과(4건)
SHOT11 = os.path.join(CACHE, "11.png")  # AI 챗봇 대화 화면
IMG_W, IMG_H = 1815, 866
IMG_AR = IMG_H / IMG_W  # 0.477

NAVY = RGBColor(0x1F, 0x2D, 0x50)
BLUE = RGBColor(0x2B, 0x5F, 0xC0)
GRAY = RGBColor(0x55, 0x5E, 0x6E)
LIGHTGRAY = RGBColor(0x9A, 0xA2, 0xB0)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x1E, 0x8A, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BD = RGBColor(0xB8, 0xC4, 0xDB)
CARD_FILL = RGBColor(0xF4, 0xF6, 0xF9)
AIBUB = RGBColor(0x60, 0x66, 0x6E)
USERBUB = RGBColor(0x2E, 0x6F, 0xD6)
CHATHDR = RGBColor(0x33, 0x37, 0x3D)
AVATAR = RGBColor(0x3D, 0x8B, 0x9E)
BTN_FILL = RGBColor(0xE3, 0xE6, 0xEA)
FONT = "맑은 고딕"


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def box(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    dash=False,
    anchor=MSO_ANCHOR.MIDDLE,
    align=PP_ALIGN.LEFT,
    round_=False,
    wrap=True,
    space_after=2,
):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    b = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    b.shadow.inherit = False
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(line_w)
        if dash:
            ln = b.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    tf = b.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        text, size, color, bold = ln
        r = p.add_run()
        set_run(r, text, size, color, bold)
    return b


def oval(slide, cx, cy, d, fill, text="", tsize=12, tcolor=WHITE, line=None):
    o = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d)
    )
    o.shadow.inherit = False
    if fill is None:
        o.fill.background()
    else:
        o.fill.solid()
        o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line
        o.line.width = Pt(2.5)
    if text:
        p = o.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text, tsize, tcolor, True)
    return o


def seg(slide, x1, y1, x2, y2, color, width=2.0, end_arrow=True):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if end_arrow:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def picture(slide, path, x, y, w, border):
    if os.path.exists(path):
        p = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
        p.line.color.rgb = border
        p.line.width = Pt(1.25)
        return p
    box(
        slide,
        x,
        y,
        w,
        w * IMG_AR,
        [("(이미지 없음)", 11, LIGHTGRAY, False)],
        fill=CARD_FILL,
        line=border,
        line_w=1.0,
    )
    return None


def header(slide, title, subtitle):
    box(
        slide,
        0.5,
        0.2,
        12.3,
        0.55,
        [(title, 21, NAVY, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    box(slide, 0.52, 0.78, 3.4, 0.06, [], fill=BLUE)
    box(
        slide,
        0.5,
        0.9,
        12.3,
        0.38,
        [(subtitle, 12.5, GRAY, False)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def footnote(slide, w=9.0):
    box(
        slide,
        0.4,
        6.55,
        w,
        0.3,
        [("※ 화면은 예시이며 실제 UI와 다를 수 있습니다.", 9.5, LIGHTGRAY, False)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def find_blank_layout(prs):
    best, best_n = prs.slide_layouts[0], len(prs.slide_layouts[0].placeholders)
    for lay in prs.slide_layouts:
        if len(lay.placeholders) < best_n:
            best, best_n = lay, len(lay.placeholders)
    return best


def clear_placeholders(slide):
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def build_semantic(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)
    header(
        s,
        "데이터, 이제 '말하듯' 찾으세요 — 의미 검색",
        "정확한 데이터 이름을 몰라도, 평소 쓰는 말로 검색하면 AI가 의미로 찾아줍니다.",
    )

    # Before: 전체 카탈로그 (작게)
    bx, by, bw = 0.35, 2.7, 3.9
    box(
        s,
        bx,
        2.32,
        bw,
        0.32,
        [("전체 카탈로그 (45건)", 11, GRAY, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        wrap=False,
    )
    picture(s, SHOT7, bx, by, bw, PANEL_BD)

    # 가운데: 일상어 검색어 콜아웃 + 화살표
    box(
        s,
        4.32,
        2.82,
        1.5,
        0.66,
        [("「배로 실어 나른 물량」", 10, NAVY, True), ("의미 검색", 11, BLUE, True)],
        fill=WHITE,
        line=BLUE,
        line_w=1.25,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        round_=True,
    )
    seg(s, 4.35, 3.72, 5.9, 3.72, BLUE, 2.5)

    # After: 의미검색 결과 (크게, 사용자가 만든 화면)
    ax, ay, aw = 5.95, 1.95, 6.95
    box(
        s,
        ax,
        1.55,
        aw,
        0.34,
        [("AI 의미검색 결과 — GLOVE 물류 실적 4건", 12.5, NAVY, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    picture(s, SHOT10, ax, ay, aw, BLUE)

    # 하단 한 줄 핵심 메시지
    box(
        s,
        0.4,
        5.65,
        12.5,
        0.68,
        [
            (
                "키워드엔 없는 '배·물량'이라는 말을 AI가 의미로 이해해, 글로벌 물류(GLOVE) 실적 데이터만 정확히 4건 골라냅니다.",
                12.5,
                NAVY,
                True,
            )
        ],
        fill=CARD_FILL,
        line=BLUE,
        line_w=1.0,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        round_=True,
    )
    footnote(s, 12.5)
    return s


def ai_bubble(slide, y, lines, h=0.5):
    oval(slide, 7.85, y + 0.18, 0.34, AVATAR)
    box(
        slide,
        8.15,
        y,
        3.0,
        h,
        lines,
        fill=AIBUB,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        round_=True,
        space_after=1,
    )


def user_bubble(slide, y, text, h=0.5):
    box(
        slide,
        8.55,
        y,
        3.05,
        h,
        [(text, 10, WHITE, False)],
        fill=USERBUB,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        round_=True,
    )


def impl_block(slide, x, y, w, n, title, subs):
    box(
        slide,
        x,
        y,
        0.34,
        0.34,
        [(str(n), 12, WHITE, True)],
        fill=BLUE,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        round_=True,
    )
    lines = [(title, 12, NAVY, True)] + [("• " + sb, 10.5, GRAY, False) for sb in subs]
    box(
        slide,
        x + 0.44,
        y - 0.06,
        w - 0.5,
        1.35,
        lines,
        anchor=MSO_ANCHOR.TOP,
        align=PP_ALIGN.LEFT,
        space_after=2,
    )


def build_chatbot(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)
    header(
        s,
        "물어보면 찾아주는 AI 챗봇",
        "어느 화면에서든 챗봇 버튼을 눌러, 평소 말투로 물어보면 데이터를 찾아 답해줍니다.",
    )

    # 좌측: 카탈로그(플로팅 버튼) + 챗봇 캡처
    ix, iy, iw = 0.35, 1.6, 4.4
    sc = iw / IMG_W
    picture(s, SHOT7, ix, iy, iw, PANEL_BD)
    bx, by = ix + 1740 * sc, iy + 795 * sc
    oval(s, bx, by, 0.5, None, line=RED)  # 빨간 강조 링
    box(
        s,
        0.45,
        iy + 2.18,
        4.2,
        0.32,
        [("↑ 모든 화면 우하단 '플로팅 챗봇 버튼'", 10, RED, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    # 챗봇 실제 화면(사용자가 만든 캡처)
    cx0, cy0, cw = 4.95, 1.95, 2.9
    picture(s, SHOT11, cx0, cy0, cw, PANEL_BD)
    # 버튼 → 챗 화면 화살표
    seg(s, bx + 0.28, by - 0.05, cx0 - 0.05, cy0 + 0.7, RED, 2.25)

    # 우측: 구현 방안 (제안서 12p 3.3 챗봇 개발)
    px, pw = 8.2, 4.65
    box(
        s,
        px,
        1.5,
        pw,
        0.48,
        [("구현 방안", 14, WHITE, True)],
        fill=BLUE,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        round_=True,
    )
    impl_block(
        s,
        px,
        2.25,
        pw,
        1,
        "사용자 접근성·편의성 극대화",
        [
            "모든 화면에 플로팅 버튼 배치 → 손쉽게 AI 챗봇 접근",
            "끊김 없는 검색 경험: 챗봇에서 검색결과 페이지로 바로 이동",
        ],
    )
    impl_block(
        s,
        px,
        3.7,
        pw,
        2,
        "Agent API 연계 및 질문 구체화",
        ["Agent API를 연계한 웹 기반 챗봇 구현", "더 정확한 답변을 위해 질문을 구체화하도록 유도"],
    )
    impl_block(
        s,
        px,
        5.15,
        pw,
        3,
        "챗봇 활용성 확장 (검토)",
        ["메뉴별 사용법(매뉴얼) 안내", "반복 질의에 대한 효과적 대응"],
    )
    footnote(s, 7.5)
    return s


def main():
    prs = Presentation(SRC)
    layout = find_blank_layout(prs)
    build_semantic(prs, layout)
    build_chatbot(prs, layout)
    prs.save(SRC)
    print("saved:", SRC)
    print("total slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
