"""
VODA 데이터 카탈로그 검색 — 유저 시나리오 상세 장표 2장을 기존 PPT에 '이어붙인다'.

파일 경로: scripts/add_scenario_detail_slides.py
목적:
    사용자가 손본 슬라이드 1·2는 그대로 두고, 3·4페이지에 "좌측 구성도 + 우측 세부사항"
    형식(DFOCUS 풍)으로 두 시나리오를 추가한다.
    - 슬라이드 3: 시나리오 ① VODA 카탈로그 화면 의미 검색
    - 슬라이드 4: 시나리오 ② AI 서비스 레이어 채팅창 검색
    구성도(3개 컬럼)는 두 장 공통, 강조/번호/세부사항만 다르게 한다.
주의:
    기존 파일을 열어 슬라이드만 추가하고 저장한다(.bak/.bak2 백업 별도 보관).
주요 의존성:
    python-pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = r"C:\Users\edgar\Desktop\데이터카탈로그 소개서\VODA_AI_Architecture.pptx"
SHOT = r"C:\Users\edgar\.claude\image-cache\c1de5183-685f-405c-bc75-472b2970def3\1.png"

ORANGE = RGBColor(0xE8, 0x77, 0x2A)
YELLOW = RGBColor(0xFF, 0xC0, 0x00)
ITEM = RGBColor(0xEC, 0xED, 0xEF)
ITEM_HI = RGBColor(0xF7, 0xE0, 0xCB)
RED = RGBColor(0xC0, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAYTXT = RGBColor(0x59, 0x59, 0x59)
GRAYLINE = RGBColor(0x8A, 0x8A, 0x8A)
DASHCOL = RGBColor(0x40, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

ROW_Y = [2.3, 3.25, 4.2]
ITEM_H = 0.8
COLS = [
    (
        0.55,
        2.35,
        "VODA 포탈 (Tableau)",
        [("cat", "카탈로그 검색 화면"), ("result", "검색결과 페이지"), ("dash", "대시보드")],
    ),
    (
        3.25,
        2.35,
        "AI 서비스 레이어 (Fabrix)",
        [("chat", "AI 서비스 채팅창"), ("agent", "AI Agent"), ("rag", "RAG 검색")],
    ),
    (
        5.95,
        2.40,
        "데이터 거버넌스",
        [
            ("extract", "메타 추출 · 인덱싱"),
            ("catalog", "메타 카탈로그"),
            ("vdb", "Knowledge / Vector DB"),
        ],
    ),
]
RECT = {}  # key -> (x, y, w, h)


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def box(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    dash=False,
    anchor=MSO_ANCHOR.MIDDLE,
    align=PP_ALIGN.CENTER,
    round_=False,
    wrap=True,
    space_after=2,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    b = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    b.shadow.inherit = False
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(line_w)
        if dash:
            ln = b.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    tf = b.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        text, size, color, bold = ln
        r = p.add_run()
        set_run(r, text, size, color, bold)
    return b


def circle(slide, cx, cy, num):
    d = 0.34
    o = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d)
    )
    o.shadow.inherit = False
    o.fill.solid()
    o.fill.fore_color.rgb = YELLOW
    o.line.color.rgb = DARK
    o.line.width = Pt(1.0)
    tf = o.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, str(num), 11, DARK, True)


def seg(slide, x1, y1, x2, y2, color, width=1.75, end_arrow=False, dash=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if end_arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def arrow(slide, pts, color, label="", lx=0, ly=0, lw=0.9, dash=False):
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        seg(slide, x1, y1, x2, y2, color, end_arrow=(i == len(pts) - 2), dash=dash)
    if label:
        box(
            slide,
            lx - lw / 2,
            ly - 0.14,
            lw,
            0.28,
            [(label, 8.5, DARK, False)],
            fill=ITEM_HI,
            line=color,
            line_w=0.75,
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
            wrap=False,
        )


def right(k):
    x, y, w, h = RECT[k]
    return (x + w, y + h / 2)


def left(k):
    x, y, w, h = RECT[k]
    return (x, y + h / 2)


def top(k):
    x, y, w, h = RECT[k]
    return (x + w / 2, y)


def bottom(k):
    x, y, w, h = RECT[k]
    return (x + w / 2, y + h)


def cyk(k):
    x, y, w, h = RECT[k]
    return y + h / 2


def draw_skeleton(slide, title):
    # 제목 + 오렌지 밑줄
    box(
        slide,
        0.5,
        0.14,
        12.0,
        0.44,
        [(title, 17, DARK, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    box(slide, 0.5, 0.64, 12.33, 0.07, [], fill=ORANGE)
    # 섹션 헤더
    box(
        slide,
        0.5,
        0.9,
        7.85,
        0.42,
        [("구성도", 13, WHITE, True)],
        fill=ORANGE,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
    )
    box(
        slide,
        8.5,
        0.9,
        4.35,
        0.42,
        [("세부사항", 13, WHITE, True)],
        fill=ORANGE,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
    )
    # 컬럼 컨테이너 + 아이템
    for cx, cw, header, items in COLS:
        box(slide, cx, 1.5, cw, 5.1, [], line=DASHCOL, line_w=1.5, dash=True, round_=True)
        box(
            slide,
            cx,
            1.6,
            cw,
            0.4,
            [(header, 11.5, DARK, True)],
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
            wrap=False,
        )
        for i, (key, label) in enumerate(items):
            RECT[key] = (cx + 0.15, ROW_Y[i], cw - 0.3, ITEM_H)


def draw_items(slide, highlight):
    for cx, cw, header, items in COLS:
        for key, label in items:
            x, y, w, h = RECT[key]
            fill = ITEM_HI if key in highlight else ITEM
            box(
                slide,
                x,
                y,
                w,
                h,
                [(label, 11, DARK, key in highlight)],
                fill=fill,
                anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.CENTER,
            )


def detail_panel(slide, blocks):
    box(slide, 8.5, 1.5, 4.35, 5.1, [], line=DASHCOL, line_w=1.5, dash=True, round_=True)
    lines = []
    for n, (title, subs) in enumerate(blocks, 1):
        lines.append((f"{n}. {title}", 11.5, DARK, True))
        for sub in subs:
            lines.append((f"   • {sub}", 10.5, GRAYTXT, False))
        lines.append(("", 6, DARK, False))
    box(
        slide,
        8.7,
        1.75,
        3.95,
        4.7,
        lines,
        anchor=MSO_ANCHOR.TOP,
        align=PP_ALIGN.LEFT,
        space_after=3,
    )


def add_shot(slide, x, y, w, label, use_image=True):
    if use_image and os.path.exists(SHOT):
        pic = slide.shapes.add_picture(SHOT, Inches(x), Inches(y), width=Inches(w))
        pic.line.color.rgb = RED
        pic.line.width = Pt(1.5)
    else:
        box(
            slide,
            x,
            y,
            w,
            w * 0.6,
            [(label, 10, GRAYTXT, False)],
            fill=ITEM,
            line=RED,
            line_w=1.0,
        )


def footnote_shot(slide):
    box(
        slide,
        8.5,
        6.65,
        4.35,
        0.3,
        [("※ 본 데이터는 가상의 예시 데이터입니다.", 9, RED, False)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def draw_governance(slide):
    """공통 거버넌스·인덱싱 흐름(회색, 배경)."""
    arrow(slide, [bottom("extract"), top("catalog")], GRAYLINE)
    arrow(
        slide,
        [bottom("catalog"), top("vdb")],
        GRAYLINE,
        "인덱싱",
        RECT["vdb"][0] + RECT["vdb"][2] / 2,
        (RECT["catalog"][1] + RECT["catalog"][3] + RECT["vdb"][1]) / 2,
        0.8,
    )
    arrow(
        slide,
        [right("rag"), left("vdb")],
        GRAYLINE,
        "조회",
        (RECT["rag"][0] + RECT["rag"][2] + RECT["vdb"][0]) / 2,
        cyk("rag"),
        0.65,
    )


def find_blank_layout(prs):
    best, best_n = prs.slide_layouts[0], len(prs.slide_layouts[0].placeholders)
    for lay in prs.slide_layouts:
        if len(lay.placeholders) < best_n:
            best, best_n = lay, len(lay.placeholders)
    return best


def clear_placeholders(slide):
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def build_scenario1(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)
    draw_skeleton(s, "사용자 시나리오 ① — VODA 카탈로그 화면 의미(Semantic) 검색")
    draw_governance(s)
    # 관계선 (빨강 = 시나리오 흐름)
    arrow(
        s,
        [right("cat"), (3.08, cyk("cat")), (3.08, cyk("agent")), left("agent")],
        RED,
        "질의",
        3.08,
        (cyk("cat") + cyk("agent")) / 2,
        0.7,
    )
    arrow(
        s,
        [
            (RECT["agent"][0], cyk("agent") + 0.18),
            (3.08, cyk("agent") + 0.18),
            (3.08, cyk("result") + 0.1),
            (RECT["result"][0] + RECT["result"][2], cyk("result") + 0.1),
        ],
        RED,
        "결과 반환",
        3.08,
        cyk("result") + 0.1,
        1.05,
    )
    arrow(
        s,
        [bottom("agent"), top("rag")],
        RED,
        "RAG 검색",
        RECT["rag"][0] + RECT["rag"][2] / 2,
        (RECT["agent"][1] + RECT["agent"][3] + RECT["rag"][1]) / 2,
        1.1,
    )
    draw_items(s, {"cat", "agent", "result"})
    circle(s, RECT["cat"][0] + 0.16, RECT["cat"][1] + 0.16, 1)
    circle(s, RECT["agent"][0] + 0.16, RECT["agent"][1] + 0.16, 2)
    circle(s, RECT["result"][0] + 0.16, RECT["result"][1] + 0.16, 3)
    detail_panel(
        s,
        [
            (
                "검색 질의",
                ["카탈로그 검색 화면에서 자연어·키워드 질의를", "AI Agent(Fabrix)에 전달"],
            ),
            (
                "의미 검색 (RAG)",
                [
                    "Hybrid(Lexical BM25+Vector)+ReRank+동의어",
                    "사전으로 메타 카탈로그 검색",
                    "권한·보안등급 필터 적용(부문별 차별 검색)",
                ],
            ),
            (
                "결과 반환",
                [
                    "유사도순 결과를 검색결과 페이지(표·필터)에 표시",
                    "데이터 원본 선택 시 Tableau 리포트 생성",
                ],
            ),
        ],
    )
    add_shot(s, 4.4, 5.15, 2.6, "예시 화면", use_image=True)
    footnote_shot(s)
    return s


def build_scenario2(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)
    draw_skeleton(s, "사용자 시나리오 ② — AI 서비스 레이어 채팅창 검색")
    draw_governance(s)
    # 관계선 (빨강)
    arrow(
        s,
        [bottom("chat"), top("agent")],
        RED,
        "질의",
        RECT["chat"][0] + RECT["chat"][2] / 2,
        (RECT["chat"][1] + RECT["chat"][3] + RECT["agent"][1]) / 2,
        0.7,
    )
    arrow(
        s,
        [bottom("agent"), top("rag")],
        RED,
        "RAG 검색",
        RECT["rag"][0] + RECT["rag"][2] / 2,
        (RECT["agent"][1] + RECT["agent"][3] + RECT["rag"][1]) / 2,
        1.1,
    )
    # 답변 반환: agent 오른쪽으로 빠져 위로 → chat 우측
    rx = RECT["agent"][0] + RECT["agent"][2]
    arrow(
        s,
        [
            (rx, cyk("agent")),
            (rx + 0.22, cyk("agent")),
            (rx + 0.22, cyk("chat")),
            (rx, cyk("chat")),
        ],
        RED,
        "답변 반환",
        rx + 0.22,
        (cyk("agent") + cyk("chat")) / 2,
        1.05,
    )
    draw_items(s, {"chat", "agent"})
    circle(s, RECT["chat"][0] + 0.16, RECT["chat"][1] + 0.16, 1)
    circle(s, RECT["agent"][0] + 0.16, RECT["agent"][1] + 0.16, 2)
    circle(s, rx + 0.22, cyk("chat") - 0.2, 3)
    detail_panel(
        s,
        [
            (
                "자연어 질의",
                ["사용자가 AI 서비스 레이어 채팅창에 자연어로 질의", "(VODA 포탈 미경유)"],
            ),
            (
                "질의 해석 · 검색 (RAG)",
                [
                    "AI Agent(Multi-Agent)가 질의 해석·분류",
                    "→ Hybrid RAG 검색",
                    "정보 부족 시 질문 구체화 유도",
                ],
            ),
            (
                "답변 생성",
                ["LLM(On-Prem)이 결과 요약 + 근거(citation) 생성", "→ 채팅창에 대화형 답변 반환"],
            ),
        ],
    )
    add_shot(s, 4.4, 5.15, 2.6, "AI 채팅 예시 화면\n(캡처 삽입)", use_image=False)
    return s


def main():
    prs = Presentation(SRC)
    layout = find_blank_layout(prs)
    build_scenario1(prs, layout)
    RECT.clear()
    build_scenario2(prs, layout)
    prs.save(SRC)
    print("saved:", SRC)
    print("total slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
