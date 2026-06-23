"""
외부 고객용(소개) 시나리오 장표 2장을 기존 PPT에 이어붙인다.

파일 경로: scripts/add_customer_slides.py
목적:
    기존 1~4페이지는 그대로 두고, 5·6페이지에 '고객 소개용' 형식(좌측 화면 / 우측 혜택)을 추가한다.
    제안서 12페이지(3.3 챗봇 개발) 톤을 참고하되, 기술 용어를 빼고 혜택 중심으로 작성한다.
    - 슬라이드 5: 카탈로그 의미(Semantic) 검색 소개
    - 슬라이드 6: AI 챗봇 소개
    좌측에는 고객이 실제 화면 캡처를 붙여 넣을 '프레임(가이드 포함)'을 둔다.
주의:
    기존 슬라이드는 건드리지 않고 추가만 한다(.bak4 백업 별도).
주요 의존성:
    python-pptx
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = r"C:\Users\edgar\Desktop\데이터카탈로그 소개서\VODA_AI_Architecture.pptx"

NAVY = RGBColor(0x1F, 0x2D, 0x50)
BLUE = RGBColor(0x2B, 0x5F, 0xC0)
BLUE_LT = RGBColor(0xEA, 0xF1, 0xFB)
FRAME_FILL = RGBColor(0xF4, 0xF6, 0xF9)
FRAME_BD = RGBColor(0xB8, 0xC4, 0xDB)
GRAY = RGBColor(0x55, 0x5E, 0x6E)
LIGHTGRAY = RGBColor(0x9A, 0xA2, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def box(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    dash=False,
    anchor=MSO_ANCHOR.MIDDLE,
    align=PP_ALIGN.LEFT,
    round_=False,
    wrap=True,
    space_after=2,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    b = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    b.shadow.inherit = False
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(line_w)
        if dash:
            ln = b.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    tf = b.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        text, size, color, bold = ln
        r = p.add_run()
        set_run(r, text, size, color, bold)
    return b


def badge(slide, cx, cy, num):
    d = 0.46
    o = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d)
    )
    o.shadow.inherit = False
    o.fill.solid()
    o.fill.fore_color.rgb = BLUE
    o.line.fill.background()
    tf = o.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, str(num), 14, WHITE, True)


def header_bar(slide, title, subtitle):
    box(
        slide,
        0.5,
        0.22,
        12.3,
        0.56,
        [(title, 21, NAVY, True)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )
    box(slide, 0.52, 0.8, 3.4, 0.06, [], fill=BLUE)
    box(
        slide,
        0.5,
        0.92,
        12.3,
        0.4,
        [(subtitle, 13, GRAY, False)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def left_frame(slide, guide_lines, floating=False):
    box(
        slide,
        0.5,
        1.5,
        7.0,
        5.0,
        guide_lines,
        fill=FRAME_FILL,
        line=FRAME_BD,
        line_w=1.5,
        dash=True,
        round_=True,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        space_after=4,
    )
    if floating:
        # 우하단 플로팅 챗봇 버튼 힌트
        d = 0.7
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.55), Inches(5.5), Inches(d), Inches(d))
        o.shadow.inherit = False
        o.fill.solid()
        o.fill.fore_color.rgb = BLUE
        o.line.color.rgb = WHITE
        o.line.width = Pt(1.5)
        tf = o.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, "AI", 13, WHITE, True)


def benefit_panel(slide, header, blocks):
    box(
        slide,
        7.85,
        1.5,
        5.0,
        0.52,
        [(header, 15, WHITE, True)],
        fill=BLUE,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
        round_=True,
    )
    y = 2.35
    for i, (title, desc) in enumerate(blocks, 1):
        badge(slide, 8.2, y + 0.23, i)
        box(
            slide,
            8.6,
            y - 0.02,
            4.2,
            0.42,
            [(title, 13.5, NAVY, True)],
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.LEFT,
            wrap=True,
        )
        box(
            slide,
            8.6,
            y + 0.42,
            4.2,
            0.85,
            [(desc, 11.5, GRAY, False)],
            anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT,
            wrap=True,
        )
        y += 1.42


def footnote(slide):
    box(
        slide,
        0.5,
        6.62,
        7.0,
        0.3,
        [("※ 화면은 예시이며 실제 UI와 다를 수 있습니다.", 9.5, LIGHTGRAY, False)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.LEFT,
        wrap=False,
    )


def find_blank_layout(prs):
    best, best_n = prs.slide_layouts[0], len(prs.slide_layouts[0].placeholders)
    for lay in prs.slide_layouts:
        if len(lay.placeholders) < best_n:
            best, best_n = lay, len(lay.placeholders)
    return best


def clear_placeholders(slide):
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def build_catalog(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)
    header_bar(
        s,
        "데이터 카탈로그를 '말하듯' 검색 — 의미(Semantic) 검색",
        "정확한 명칭을 몰라도, 찾고 싶은 데이터를 자연어로 입력하면 AI가 의미로 찾아드립니다.",
    )
    left_frame(
        s,
        [
            ("[ 화면 캡처 삽입 위치 ]", 14, LIGHTGRAY, True),
            ("카탈로그 의미검색 화면 캡처를 여기에 붙여 넣으세요", 11, LIGHTGRAY, False),
        ],
    )
    benefit_panel(
        s,
        "이렇게 달라집니다",
        [
            (
                "키워드를 몰라도 검색돼요",
                '"해상 운송 월별 실적"처럼 일상어로 입력해도 의미가 통하는 데이터를 찾아줍니다.',
            ),
            (
                "결과가 더 정확해요",
                "키워드와 의미를 함께 보고 정렬해 원하는 데이터를 빠르게 찾고, 권한에 맞는 데이터만 보여줍니다.",
            ),
            (
                "찾으면 바로 활용해요",
                "원하는 데이터 원본을 선택하면 곧바로 리포트 작성 화면으로 이어집니다.",
            ),
        ],
    )
    footnote(s)
    return s


def build_chatbot(prs, layout):
    s = prs.slides.add_slide(layout)
    clear_placeholders(s)
    header_bar(
        s,
        "물어보면 찾아주는 AI 데이터 도우미 — AI 챗봇",
        "어느 화면에서든 챗봇을 열고 자연어로 물어보면, 필요한 데이터를 찾아 답해드립니다.",
    )
    left_frame(
        s,
        [
            ("[ 화면 캡처 삽입 위치 ]", 14, LIGHTGRAY, True),
            ("AI 챗봇 대화 화면 캡처를 여기에 붙여 넣으세요", 11, LIGHTGRAY, False),
            ("(우하단 플로팅 챗봇 버튼 포함)", 10, LIGHTGRAY, False),
        ],
        floating=True,
    )
    benefit_panel(
        s,
        "이렇게 달라집니다",
        [
            (
                "어디서나 바로 물어봐요",
                "모든 화면 우측 하단의 플로팅 버튼으로, 보던 화면을 떠나지 않고 바로 질문합니다.",
            ),
            (
                "대화로 쉽게 좁혀가요",
                "무엇을 물어야 할지 몰라도, 챗봇이 더 좋은 질문을 하도록 안내해 원하는 답에 다가갑니다.",
            ),
            (
                "근거와 함께 답해줘요",
                "찾은 데이터의 출처를 함께 알려주고, 검색 결과 화면으로 바로 이어집니다.",
            ),
        ],
    )
    footnote(s)
    return s


def main():
    prs = Presentation(SRC)
    layout = find_blank_layout(prs)
    build_catalog(prs, layout)
    build_chatbot(prs, layout)
    prs.save(SRC)
    print("saved:", SRC)
    print("total slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
