"""Generate SERVICE_OVERVIEW.pptx from docs/SERVICE_OVERVIEW.md outline.

Run: python scripts/generate_overview_ppt.py
Output: docs/SERVICE_OVERVIEW.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

FONT_KR = "맑은 고딕"
FONT_MONO = "Consolas"

COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)      # 진청
COLOR_ACCENT = RGBColor(0x2E, 0x7D, 0x32)       # 진녹
COLOR_WARN = RGBColor(0xE6, 0x5C, 0x00)         # 진주황
COLOR_TEXT = RGBColor(0x21, 0x21, 0x21)
COLOR_SUBTLE = RGBColor(0x5F, 0x6B, 0x7A)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_BG_BOX = RGBColor(0xE8, 0xEE, 0xF7)
COLOR_BORDER = RGBColor(0xC4, 0xCE, 0xDB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_font(run, *, size=14, bold=False, color=None, mono=False):
    run.font.name = FONT_MONO if mono else FONT_KR
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=None,
              align=PP_ALIGN.LEFT, mono=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color or COLOR_TEXT, mono=mono)
    return tb


def _add_box(slide, left, top, width, height, *, fill=COLOR_BG_BOX, border=COLOR_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _add_header_bar(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False

    _add_text(slide, Inches(0.5), Inches(0.1), Inches(10), Inches(0.5), title,
              size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        _add_text(slide, Inches(0.5), Inches(0.55), Inches(10), Inches(0.3), subtitle,
                  size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    page_num_tb = slide.shapes.add_textbox(Inches(12.2), Inches(0.25), Inches(1.0), Inches(0.4))
    page_num_tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return bar


def _add_bullet_list(slide, left, top, width, height, items: list[str], *, size=14):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        _set_font(run, size=size, color=COLOR_TEXT)
    return tb


def _add_table(slide, left, top, width, height, headers: list[str], rows: list[list[str]],
               *, header_size=11, body_size=10, col_widths: list[float] | None = None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    if col_widths is not None:
        total = sum(col_widths)
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = Emu(int(width * (w / total)))

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        _set_font(run, size=header_size, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_LIGHT if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            _set_font(run, size=body_size, color=COLOR_TEXT)
    return table


def _slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃


# ── 슬라이드 빌더 ────────────────────────────────────────────────────────


def slide_cover(prs):
    s = _slide_blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    bg.shadow.inherit = False

    _add_text(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
              "보험청구심사 어시스턴트",
              size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(3.4), Inches(11.3), Inches(0.6),
              "서비스 종합 안내서",
              size=22, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(4.4), Inches(11.3), Inches(0.5),
              "RAG 기반 대국민 보험 청구 가능성 사전 확인 서비스",
              size=14, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    line.line.fill.background()
    line.shadow.inherit = False

    _add_text(s, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
              "작성일: 2026-05-24    |    대상: 외부 공유 · 소개 · 온보딩",
              size=13, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.4),
              "테스트: 898 passed · ruff 0 · Sprint 8.6 완료",
              size=12, color=RGBColor(0x9C, 0xAA, 0xC0), align=PP_ALIGN.CENTER)
    return s


def slide_toc(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "목차", "Service Overview — 11 sections")
    items = [
        ("1", "서비스 개요", "도메인 · 사용자 · 진행률"),
        ("2", "기술 스택", "백엔드 · AI · DB · 프론트 · 인프라"),
        ("3", "전체 아키텍처", "레이어별 데이터 흐름"),
        ("4", "API 명세 요약", "세션 4 · 문서 2 · 운영 2"),
        ("5", "프론트엔드 컴포넌트", "React + Vite + TypeScript"),
        ("6", "뉴로심볼릭 구성", "Neuro + Symbolic 결합 (~93%)"),
        ("7", "시나리오 흐름 4건", "자동차 · 화재 × 정보충분 · 부족"),
        ("8", "데이터 자산", "PDF 4 · 청크 739 · 그래프 748 노드"),
        ("9", "운영 인프라", "SLO · 감사 · PII · rate limit"),
        ("10", "로드맵", "Sprint 9~12+"),
        ("11", "참고 문서 인덱스", "design / requirements / usage / pm"),
    ]
    top = Inches(1.2)
    row_h = Inches(0.5)
    for idx, (num, title, desc) in enumerate(items):
        y = top + row_h * idx
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.05),
                                     Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_PRIMARY
        circle.line.fill.background()
        circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        _set_font(r, size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        _add_text(s, Inches(1.6), y, Inches(3.5), row_h, title,
                  size=15, bold=True, color=COLOR_PRIMARY)
        _add_text(s, Inches(5.2), y + Inches(0.05), Inches(7.5), row_h, desc,
                  size=12, color=COLOR_SUBTLE)
    return s


def slide_overview(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "1. 서비스 개요", "한 줄 소개 + 도메인 + 사용자")

    _add_box(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.4),
             fill=COLOR_BG_BOX)
    _add_text(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(1.2),
              "일반 국민이 자신의 보험금 청구 가능성을 자연어로 질문하면, 실제 보험사 약관 PDF 원문을 근거로 "
              "'가능성 높음 / 중간 / 낮음' 등급과 조항 인용을 함께 제시하는 RAG 기반 대국민 서비스 도구.\n\n"
              "단정적 판단('된다 / 안된다') 대신 충족 항목과 미충족 항목을 정리해 사용자가 직접 최종 판단을 "
              "내릴 수 있도록 돕는다. 모든 응답에는 법적 면책 문구가 포함된다.",
              size=13, color=COLOR_TEXT)

    _add_text(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.4),
              "대상 영역", size=14, bold=True, color=COLOR_ACCENT)
    _add_bullet_list(s, Inches(0.8), Inches(3.1), Inches(5.8), Inches(1.5), [
        "자동차보험 (auto)",
        "화재보험 (fire)",
        "상해질병 (accident_disease)",
        "현재 데이터: 한화손해보험 자동차·화재 약관 + 요약서 PDF 4종",
    ], size=12)

    _add_text(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.4),
              "주요 사용자", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(6.8), Inches(3.1), Inches(6.2), Inches(1.6),
               ["사용자", "진입 방식", "목적"],
               [
                   ["일반 국민", "웹 UI 채팅", "청구 가능성 사전 확인"],
                   ["운영자", "CLI + API + /metrics", "데이터 적재 · 모니터링 · 감사"],
                   ["개발자", "소스코드 + 본 문서", "온보딩 · 기능 이해"],
               ],
               col_widths=[1.2, 2.4, 2.6])

    _add_text(s, Inches(0.5), Inches(5.0), Inches(6), Inches(0.4),
              "단정적 판단 회피 원칙", size=14, bold=True, color=COLOR_WARN)
    _add_bullet_list(s, Inches(0.8), Inches(5.4), Inches(12), Inches(1.8), [
        "satisfied / unsatisfied 두 목록으로 분리해서 제시 — 사용자가 직접 판단",
        "likelihood 는 '높음 / 중간 / 낮음' 3단계 + confidence (full / partial) 명시",
        "partial: 필수 슬롯 일부 미충족 상태에서의 추정 응답 — UI '(추정)' 배지",
        "모든 assessment 응답에 면책 문구 의무 포함",
    ], size=12)


def slide_progress(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "1. 서비스 개요 — 진행률", "2026-05-24 기준 · 898 tests · ruff 0")

    rows = [
        ["1", "PDF 적재 파이프라인 + CLI", "완료"],
        ["2", "멀티턴 대화 HTTP API", "완료"],
        ["3", "웹 UI 사양서 + 백엔드 정비", "완료"],
        ["4", "GraphRAG + Hybrid + ReAct 골격", "완료"],
        ["5", "인용 카드 PDF 페이지 캡처", "완료"],
        ["6", "응답 품질 정책 (partial + 모름 처리)", "완료"],
        ["7", "응답 톤 정책 (능동적 안내 + 친절체)", "완료"],
        ["8", "대국민 서비스 전환 인프라 (감사/PII/rate/circuit)", "완료"],
        ["8.5", "후속 보정 + 디자인 패키지 + frontend 통합", "완료"],
        ["8.6", "옵션 노출 정책 정교화 (OptionsPanel)", "완료"],
        ["9", "외부 read-only tool (KIDI 활성 + law/hira 골격)", "진행 ~60%"],
        ["10", "fss 금감원 크롤링", "보류"],
        ["11", "ReAct agent 본격 활성 (AgentRunner)", "진행 ~93%"],
        ["12+", "외부 호스팅 · 도메인 · DPIA", "미착수"],
    ]
    _add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8),
               ["Sprint", "내용", "상태"], rows,
               col_widths=[0.8, 6.5, 1.8], header_size=12, body_size=11)


def slide_stack_backend(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "2. 기술 스택 — 백엔드", "Python 3.11+ · FastAPI · SQLAlchemy")

    _add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.5),
               ["기술", "역할", "선택 이유"],
               [
                   ["Python 3.11+", "전체 백엔드 언어", "AI/ML 라이브러리 생태계 + 팀 역량"],
                   ["FastAPI 0.110+", "REST API 서버", "비동기 · 타입 검증 · OpenAPI 자동 문서"],
                   ["Typer + Rich", "CLI (ica 명령어)", "데이터 파이프라인 운영자 도구"],
                   ["SQLAlchemy + Alembic", "ORM + 마이그레이션", "SQL injection 차단 + 스키마 이력"],
                   ["pydantic-settings", "환경 변수 관리", ".env 기반 타입 안전 설정"],
                   ["slowapi", "rate limit (per-IP / per-session)", "FastAPI 네이티브 미들웨어"],
                   ["pybreaker", "circuit breaker", "외부 API 장애 격리"],
                   ["httpx", "외부 API 비동기 클라이언트", "MCP 없이 직접 REST 호출"],
               ],
               col_widths=[2.5, 4.5, 5.3], header_size=12, body_size=11)


def slide_stack_ai(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "2. 기술 스택 — AI / 검색", "gpt-4o-mini · Chroma · Neo4j · LangChain")

    _add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.0),
               ["기술", "역할", "선택 이유"],
               [
                   ["OpenAI gpt-4o-mini", "LLM (슬롯 추출 · 질문 생성 · 판단)",
                    "한국어 품질 무난 + 비용 저렴"],
                   ["text-embedding-3-small (1536-d)", "청크 임베딩",
                    "LLM과 동일 제공자 → API 키 1개"],
                   ["Chroma (로컬 영속화)", "벡터 DB",
                    "임베디드 (별도 서버 없음) + 메타 필터링"],
                   ["Neo4j 5.x community (Docker)", "그래프 DB",
                    "약관 계층 탐색 + LangChain GraphCypherQAChain"],
                   ["LangChain (app/rag/ 한정)", "GraphRAG 래퍼",
                    "GraphCypherQAChain 재사용 + LangChain-neo4j 통합"],
                   ["pdfplumber + PyMuPDF", "PDF 파서",
                    "표 추출 + 페이지 이미지 (Sprint 5 lazy 변환)"],
               ],
               col_widths=[3.2, 4.2, 4.9], header_size=12, body_size=11)

    _add_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), fill=COLOR_BG_BOX)
    _add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4),
              "LLM 호출 3종 (모두 Function Calling / Structured Output)",
              size=13, bold=True, color=COLOR_PRIMARY)
    _add_text(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.1),
              "extract_slots (temperature 0.0, 매 turn 의무)   ·   "
              "next_question (temperature 0.0, ask 응답)   ·   "
              "generate_assessment (temperature 0.2, JSON Schema 강제)",
              size=12, color=COLOR_TEXT)


def slide_stack_data_infra(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "2. 기술 스택 — DB · 프론트 · 인프라", "SQLite/PostgreSQL · React · Docker · Prometheus")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "데이터베이스", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(2.6),
               ["DB", "용도"],
               [
                   ["SQLite / PostgreSQL", "메타데이터 + 감사 로그"],
                   ["Chroma", "벡터 임베딩 + 유사도 검색"],
                   ["Neo4j", "약관 계층 지식 그래프"],
                   ["인메모리 dict + TTL", "대화 세션 (30분)"],
               ],
               col_widths=[2.4, 3.6], header_size=11, body_size=10)

    _add_text(s, Inches(6.8), Inches(1.05), Inches(6), Inches(0.4),
              "프론트엔드", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(1.6),
               ["기술", "역할"],
               [
                   ["React + Vite + TypeScript", "채팅 UI (Claude 디자인 생성)"],
                   ["CSS Modules", "컴포넌트 스코프 스타일 격리"],
               ],
               col_widths=[2.6, 3.4], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(4.4), Inches(6), Inches(0.4),
              "운영 / 인프라", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
               ["기술", "역할"],
               [
                   ["Docker Compose", "Neo4j · PostgreSQL 로컬 실행"],
                   ["Prometheus /metrics", "SLO 메트릭 수집"],
                   ["Alembic", "DB 스키마 마이그레이션"],
                   ["presidio (옵션) · 정규식 5종", "PII 마스킹"],
                   ["cachetools", "외부 API 인메모리 캐시 (TTL 24h)"],
                   ["GitHub Actions", "CI (PR마다 pytest + ruff)"],
               ],
               col_widths=[3.2, 9.1], header_size=11, body_size=10)


def slide_architecture(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "3. 전체 아키텍처", "사용자 → 미들웨어 → LLM / RAG / Tool → 저장소 · 외부 API")

    def layer_box(left, top, width, height, title, items, color):
        outer = _add_box(s, left, top, width, height, fill=color, border=COLOR_BORDER)
        _add_text(s, left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.35),
                  title, size=12, bold=True, color=COLOR_PRIMARY)
        tb = s.shapes.add_textbox(left + Inches(0.1), top + Inches(0.4),
                                   width - Inches(0.2), height - Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(2)
            r = p.add_run(); r.text = f"• {it}"
            _set_font(r, size=10, color=COLOR_TEXT)
        return outer

    # row 1 — 클라이언트
    layer_box(Inches(0.5), Inches(1.1), Inches(6.0), Inches(1.0),
              "[1] 클라이언트", ["사용자 브라우저 (React + Vite)", "운영자 CLI (ica 명령어)"],
              RGBColor(0xE3, 0xF2, 0xFD))

    # 운영 / 외부 API 영역 (옆)
    layer_box(Inches(6.8), Inches(1.1), Inches(6.0), Inches(1.0),
              "[7] 외부 API (Sprint 9~10)",
              ["법령정보센터 / HIRA / KIDI(✅) / 금감원"],
              RGBColor(0xFF, 0xF3, 0xE0))

    # row 2 — FastAPI
    layer_box(Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.1),
              "[2] FastAPI 서버",
              ["slowapi rate limit  →  Audit Middleware (response_id)  →  PII 마스킹 (입력)  →  sessions.service.post_message"],
              RGBColor(0xE8, 0xF5, 0xE9))

    # row 3 — LLM 레이어
    layer_box(Inches(0.5), Inches(3.6), Inches(6.0), Inches(1.0),
              "[3] LLM 레이어 (gpt-4o-mini)",
              ["extract_slots → next_question → generate_assessment"],
              RGBColor(0xF3, 0xE5, 0xF5))

    # row 3 — RAG / Tool
    layer_box(Inches(6.8), Inches(3.6), Inches(6.0), Inches(1.0),
              "[4] RAG / Tool 레이어",
              ["rag.service.retrieve (Vector + Graph + Hybrid)  ·  AgentRunner (ReAct loop, Sprint 11)  ·  Tool Dispatcher (8종)"],
              RGBColor(0xFF, 0xEB, 0xEE))

    # row 4 — 저장소
    layer_box(Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.0),
              "[5] 저장소",
              ["Chroma (739 벡터)   ·   Neo4j (748 노드)   ·   SQLite / PostgreSQL (메타 + audit_log)"],
              RGBColor(0xEC, 0xEF, 0xF1))

    # row 5 — 출력
    layer_box(Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
              "[6] 응답 후처리",
              ["PII 마스킹 (출력)   →   Audit Write (llm_calls + retrieved_chunk_ids)   →   클라이언트 응답"],
              RGBColor(0xE8, 0xF5, 0xE9))


def slide_architecture_table(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "3. 전체 아키텍처 — 레이어 역할 요약", "각 레이어의 책임과 핵심 파일")

    _add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.5),
               ["레이어", "역할", "핵심 파일"],
               [
                   ["FastAPI Router", "인증 없음 + rate limit + CORS", "app/sessions/router.py"],
                   ["Audit Middleware", "response_id 생성 + 감사 기록", "app/audit/middleware.py"],
                   ["PII 마스킹", "입출력 개인정보 정규식 차단", "app/security/pii.py"],
                   ["sessions.service", "오케스트레이션 (슬롯 수집 → 분기 → 응답)", "app/sessions/service.py"],
                   ["LLM 레이어", "3종 Function Calling + Structured Output", "app/sessions/llm.py"],
                   ["RAG 레이어", "Vector / Graph / Hybrid 검색 + ReAct loop", "app/rag/"],
                   ["Tool Dispatcher", "8종 tool 라우팅", "app/tools/dispatcher.py"],
                   ["외부 어댑터", "httpx + cachetools + circuit breaker", "app/external/"],
               ],
               col_widths=[2.5, 5.8, 4.0], header_size=12, body_size=11)


def slide_api_endpoints(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "4. API 명세 요약 — 엔드포인트",
                    "Base: http://localhost:8000/api/v1   ·   인증 없음 (대국민)")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "세션 API (4종)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.4),
               ["메서드", "경로", "응답", "설명"],
               [
                   ["POST", "/sessions", "201", "새 대화 세션 생성"],
                   ["POST", "/sessions/{id}/messages", "200 ask | assessment", "멀티턴 핵심 — 슬롯 충족 분기"],
                   ["GET", "/sessions/{id}", "200", "디버그 — 슬롯 + 대화 이력"],
                   ["DELETE", "/sessions/{id}", "204", "세션 명시 폐기 (멱등)"],
               ],
               col_widths=[1.0, 3.5, 2.5, 5.3], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(4.2), Inches(6), Inches(0.4),
              "문서 API (2종)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(4.6), Inches(6.0), Inches(1.4),
               ["경로", "설명"],
               [
                   ["GET /documents/products", "등록 상품 목록"],
                   ["GET /documents/insurers", "등록 보험사 목록"],
               ],
               col_widths=[2.8, 3.2], header_size=11, body_size=10)

    _add_text(s, Inches(6.8), Inches(4.2), Inches(6), Inches(0.4),
              "운영 API (2종)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(6.8), Inches(4.6), Inches(6.0), Inches(1.4),
               ["경로", "설명"],
               [
                   ["GET /health", "서버 상태 확인"],
                   ["GET /metrics", "Prometheus 메트릭 (SLO)"],
               ],
               col_widths=[2.6, 3.4], header_size=11, body_size=10)


def slide_api_responses(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "4. API — 응답 모드 + 에러 코드", "discriminated union (type: ask | assessment)")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "type: ask — 슬롯 보강 질의", size=13, bold=True, color=COLOR_PRIMARY)
    _add_box(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.5),
             fill=COLOR_BG_LIGHT, border=COLOR_BORDER)
    ask_json = ('{\n'
                '  "assistant": {\n'
                '    "type": "ask",\n'
                '    "message": "사고 당시 과실 비율을 알고 계신가요?",\n'
                '    "expected_slots": ["fault_ratio"],\n'
                '    "options": ["0%", "10%", "20~50%", "50%+", "모르겠습니다"]\n'
                '  }\n'
                '}')
    _add_text(s, Inches(0.65), Inches(1.6), Inches(5.7), Inches(2.3),
              ask_json, size=10, mono=True, color=COLOR_TEXT)

    _add_text(s, Inches(6.8), Inches(1.05), Inches(6), Inches(0.4),
              "type: assessment — 최종 판단", size=13, bold=True, color=COLOR_PRIMARY)
    _add_box(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5),
             fill=COLOR_BG_LIGHT, border=COLOR_BORDER)
    ass_json = ('{\n'
                '  "assistant": {\n'
                '    "type": "assessment",\n'
                '    "likelihood": "중간",\n'
                '    "confidence": "full" | "partial",\n'
                '    "summary": "...",\n'
                '    "satisfied": [...], "unsatisfied": [...],\n'
                '    "citations": [...], "next_steps": [...],\n'
                '    "disclaimer": "본 결과는 참고용이며 ..."\n'
                '  }\n'
                '}')
    _add_text(s, Inches(6.95), Inches(1.6), Inches(5.7), Inches(2.3),
              ass_json, size=10, mono=True, color=COLOR_TEXT)

    _add_text(s, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
              "공통 에러 코드", size=14, bold=True, color=COLOR_WARN)
    _add_table(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.2),
               ["HTTP", "code", "발생 시점"],
               [
                   ["400", "VALIDATION_ERROR", "입력값 검증 실패"],
                   ["404", "SESSION_NOT_FOUND", "세션 만료 또는 오타"],
                   ["429", "RATE_LIMITED", "per-IP 10 req/min 초과"],
                   ["503", "LLM_UNAVAILABLE", "OpenAI 호출 실패"],
               ],
               col_widths=[1.0, 3.0, 8.3], header_size=11, body_size=11)


def slide_frontend_tree(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "5. 프론트엔드 컴포넌트 트리", "React + Vite + TypeScript · CSS Modules")

    tree = (
        "App\n"
        "├─ ChatHeader\n"
        "│   └─ NewChatButton\n"
        "├─ MessageList\n"
        "│   └─ MessageBubble\n"
        "│       ├─ AskCard\n"
        "│       │   └─ OptionsPanel  ← closed-ended 슬롯만 노출\n"
        "│       └─ AssessmentCard\n"
        "│           ├─ LikelihoodBadge\n"
        "│           ├─ PartialBadge  ← '(추정)' 표시\n"
        "│           ├─ SatisfiedList / UnsatisfiedList\n"
        "│           ├─ CitationList\n"
        "│           │   └─ CitationItem\n"
        "│           │       ├─ PageImageThumb  ← /static/page_images/\n"
        "│           │       └─ PdfLink  ← #page=N\n"
        "│           ├─ NextStepsList\n"
        "│           └─ Disclaimer\n"
        "├─ ChatInput\n"
        "└─ SlotInspector  ← 디버그 (접힘)\n"
    )
    _add_box(s, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.7), fill=COLOR_BG_LIGHT)
    _add_text(s, Inches(0.7), Inches(1.2), Inches(7.2), Inches(5.5),
              tree, size=11, mono=True, color=COLOR_TEXT)

    _add_text(s, Inches(8.3), Inches(1.1), Inches(5), Inches(0.4),
              "라우팅", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(2.4),
               ["경로", "내용"],
               [
                   ["/", "메인 채팅 단일 페이지"],
                   ["/legal", "법적 이용약관"],
                   ["/disclaimer", "서비스 면책"],
                   ["/privacy", "개인정보 처리방침"],
                   ["/accessibility", "접근성 안내"],
                   ["/sources", "데이터 출처"],
               ],
               col_widths=[1.6, 2.9], header_size=11, body_size=10)

    _add_text(s, Inches(8.3), Inches(4.2), Inches(5), Inches(0.4),
              "useSession 흐름", size=14, bold=True, color=COLOR_ACCENT)
    _add_bullet_list(s, Inches(8.3), Inches(4.6), Inches(4.7), Inches(2.5), [
        "첫 메시지 → POST /sessions → sessionId 획득",
        "이후 → POST /sessions/{id}/messages",
        "type 분기 → AskCard or AssessmentCard 렌더",
        "오류 시 낙관적 업데이트 롤백 + 에러 표시",
    ], size=11)


def slide_neurosymbolic_neuro(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "6. 뉴로심볼릭 — Neuro (신경망)",
                    "LLM 자연어 이해 + 추론   ·   현재 진행률 ~93%")

    _add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.0),
               ["구성 요소", "역할", "모델 / 적용 범위", "상태"],
               [
                   ["extract_slots", "자연어 → SlotState 필드 추출",
                    "gpt-4o-mini (T=0.0) / 매 turn", "완료"],
                   ["next_question", "부족 슬롯 → 질문 + 옵션",
                    "gpt-4o-mini (T=0.0) / ask 응답", "완료"],
                   ["generate_assessment", "슬롯 + RAG → 가능성 + 인용",
                    "gpt-4o-mini (T=0.2) / assessment", "완료"],
                   ["RAG Hybrid 검색", "질의 → 관련 약관 청크",
                    "text-embedding-3-small (1536-d)", "완료"],
                   ["ReAct agent loop", "LLM이 tool 자가 선택 · 반복",
                    "RAG_REACT=true 시 (max_iter=5)", "~93%"],
               ],
               col_widths=[2.5, 4.0, 3.8, 2.0], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
              "LLM 호출 순서 (기본 흐름)", size=14, bold=True, color=COLOR_PRIMARY)
    flow = "extract_slots  →  (슬롯 부족) next_question  /  (충족) RAG 검색  →  generate_assessment"
    _add_box(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.9), fill=COLOR_BG_BOX)
    _add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
              flow, size=13, color=COLOR_TEXT, align=PP_ALIGN.CENTER)


def slide_neurosymbolic_symbolic(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "6. 뉴로심볼릭 — Symbolic (규칙 기반)",
                    "결정론 계산 + 신뢰성 보장")

    _add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.7),
               ["구성 요소", "역할", "위치"],
               [
                   ["슬롯 validator (_compute_missing)", "영역별 필수 슬롯 충족 여부 결정",
                    "app/sessions/service.py"],
                   ["partial 분기 (_should_partial)", "ask≥3 / unknown≥2 / 명시 키워드",
                    "app/sessions/service.py"],
                   ["Neo4j 지식 그래프", "Insurer→Product→...→Clause 탐색",
                    "app/rag/graph.py"],
                   ["Tool Dispatcher", "LLM tool_call → 함수 라우팅 (8종)",
                    "app/tools/dispatcher.py"],
                   ["KIDI 과실비율 정적 데이터", "6 시나리오 (get_fault_ratio_standard)",
                    "app/external/kidi/"],
                   ["calc_claim_amount", "의료수가/손해액 → 보험금 산정",
                    "app/tools/calc.py"],
                   ["validate_coverage_period", "사고일 ∈ 보장기간 검증",
                    "app/tools/calc.py"],
                   ["구조 인식 청킹", "제N조/항/표 단위 정규식 분할",
                    "app/chunks/"],
               ],
               col_widths=[3.5, 4.8, 4.0], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
              "옵션 노출 정책 (Sprint 8.6 — Claude Plan 모드 패턴)",
              size=13, bold=True, color=COLOR_WARN)
    _add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8),
              "closed-ended 5종 (area, incident_type, damage_type, loss_type, cause) → options 채움 + '모르겠습니다' 의무\n"
              "open-ended (insurer, incident_date, diagnosis, fault_ratio 등) → options 빈 배열 → OptionsPanel 자동 숨김",
              size=11, color=COLOR_TEXT)


def _scenario_slide(prs, num: str, title: str, subtitle: str, steps: list[tuple[str, str]],
                    result_box: str):
    """Common scenario slide layout: turn-by-turn dialogue + result."""
    s = _slide_blank(prs)
    _add_header_bar(s, f"7. 시나리오 {num} — {title}", subtitle)

    _add_text(s, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
              "대화 흐름", size=13, bold=True, color=COLOR_ACCENT)

    top = Inches(1.5)
    row_h = Inches(0.62)
    for i, (actor, msg) in enumerate(steps):
        y = top + row_h * i
        is_user = actor.startswith("사용자")
        bg = RGBColor(0xE3, 0xF2, 0xFD) if is_user else RGBColor(0xF1, 0xF8, 0xE9)
        # actor badge
        badge = _add_box(s, Inches(0.5), y + Inches(0.05), Inches(1.5), Inches(0.45),
                         fill=COLOR_PRIMARY if is_user else COLOR_ACCENT,
                         border=COLOR_BORDER)
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = actor
        _set_font(r, size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        _add_box(s, Inches(2.2), y, Inches(10.6), row_h - Inches(0.05), fill=bg)
        _add_text(s, Inches(2.35), y + Inches(0.05), Inches(10.3), row_h - Inches(0.1),
                  msg, size=10, color=COLOR_TEXT)

    _add_box(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0),
             fill=COLOR_BG_BOX, border=COLOR_PRIMARY)
    _add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.35),
              "결과", size=12, bold=True, color=COLOR_PRIMARY)
    _add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.5),
              result_box, size=11, color=COLOR_TEXT)


def slide_scenario_a(prs):
    _scenario_slide(prs, "A", "자동차 정보 충분 → ask 1턴 → assessment full",
                    "보험사·사고유형·과실비율을 첫 메시지에 포함한 경우",
                    [
                        ("사용자", "한화손해보험 자동차보험. 어제 추돌사고로 상대방 차 긁혔어요. 과실비율은 20%입니다."),
                        ("LLM", "extract_slots → area=auto, insurer=한화, incident_type=추돌, fault_ratio=20"),
                        ("백엔드", "_compute_missing → incident_date 부족 → next_question 호출"),
                        ("응답", "ask: '사고 발생일이 언제인가요?'"),
                        ("사용자", "2026-05-10이요"),
                        ("백엔드", "모든 필수 슬롯 충족 → RAG retrieve(hybrid) → top-8 chunks"),
                        ("응답", "assessment: likelihood=높음, confidence=full, citations=[제15조①, 제22조]"),
                    ],
                    "✅ assessment full · likelihood=높음 · 한화 자동차약관 인용 2건 포함")


def slide_scenario_b(prs):
    _scenario_slide(prs, "B", "자동차 정보 부족 → 즉시 partial",
                    "'그냥 모름' 시나리오 — unknown_slots ≥ 2 → partial 분기",
                    [
                        ("사용자", "어제 자동차 사고 났어요. 한화보험인데 나머지는 잘 모르겠어요"),
                        ("LLM", "extract_slots → area=auto, insurer=한화, unknown_slots=[incident_type, fault_ratio]"),
                        ("백엔드", "_should_partial → unknown ≥ 2 → partial 분기 즉시 진입"),
                        ("백엔드", "RAG retrieve(hybrid, slots={area, insurer}) → 자동차 일반 보장 chunks"),
                        ("LLM", "generate_assessment(partial=true)"),
                        ("응답", "assessment: likelihood=중간, confidence=partial, summary='정보 일부 부족 — 일반 기준 안내'"),
                    ],
                    "⚠ assessment partial · UI '(추정)' 배지 표시 · 면책 강조")


def slide_scenario_c(prs):
    _scenario_slide(prs, "C", "화재 정보 충분 → ask 2턴 → assessment full",
                    "화재 영역 필수 슬롯 — loss_type / cause 추가 수집",
                    [
                        ("사용자", "한화손해보험 화재보험인데 부엌에서 불이 났어요"),
                        ("LLM", "extract_slots → area=fire, insurer=한화, cause=가스/조리부주의"),
                        ("응답", "ask + OptionsPanel: '피해 정도?' [전손, 부분손해, 도난, 기타, 모르겠습니다]"),
                        ("사용자", "부분손해요"),
                        ("응답", "ask: '피해 입은 물품과 사고 날짜를 알려주세요' (open-ended → OptionsPanel 숨김)"),
                        ("사용자", "가전제품이요. 2026-05-15"),
                        ("백엔드", "모든 필수 슬롯 충족 → RAG retrieve → 화재보험 약관 chunks"),
                        ("응답", "assessment: likelihood=중간, confidence=full"),
                    ],
                    "✅ assessment full · closed-ended 옵션 → open-ended 자동 분기 시연")


def slide_scenario_d(prs):
    _scenario_slide(prs, "D", "화재 정보 부족 → 즉시 partial",
                    "화재 사고 정보를 거의 모르는 사용자",
                    [
                        ("사용자", "집에 화재가 났는데 보험금 받을 수 있을까요? 자세한 건 다 모르겠어요"),
                        ("LLM", "extract_slots → area=fire, unknown_slots=[insurer, loss_type, cause, damaged_items]"),
                        ("사용자", "그냥 알려주세요"),
                        ("백엔드", "ask 횟수 ≥ 1 + unknown ≥ 2 + '그냥' 키워드 → partial 강제"),
                        ("백엔드", "RAG retrieve(vector, slots={area=fire}) → 화재 일반 보장"),
                        ("LLM", "generate_assessment(partial=true)"),
                        ("응답", "assessment: confidence=partial, summary='현재 정보로 일반 기준 안내'"),
                    ],
                    "⚠ assessment partial · 3 partial 트리거 동시 충족 (ask+unknown+키워드)")


def slide_data_pdf_sql(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "8. 데이터 자산 — PDF + SQLite",
                    "한화손해보험 자동차·화재 약관 4종 · 메타데이터 6 테이블")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "PDF 원본 데이터 (4종)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.4),
               ["보험사", "영역", "문서 종류", "파일"],
               [
                   ["한화손해보험", "auto", "약관 (terms)", "data/raw/hanwha/auto/.../terms.pdf"],
                   ["한화손해보험", "auto", "상품요약 (summary)", "data/raw/hanwha/auto/.../summary.pdf"],
                   ["한화손해보험", "fire", "약관 (terms)", "data/raw/hanwha/fire/.../terms.pdf"],
                   ["한화손해보험", "fire", "상품요약 (summary)", "data/raw/hanwha/fire/.../summary.pdf"],
               ],
               col_widths=[2.2, 1.0, 2.5, 6.6], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(4.2), Inches(6), Inches(0.4),
              "SQLite / PostgreSQL 메타데이터", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5),
               ["테이블", "레코드 수", "내용"],
               [
                   ["insurers", "1", "한화손해보험"],
                   ["products", "2", "개인용자동차보험, 화재보험"],
                   ["product_versions", "2", "각 상품 현행 판매 버전"],
                   ["documents", "4", "terms + summary × 2 영역"],
                   ["clause_chunks", "739", "구조 인식 청킹 결과"],
                   ["audit_log", "17+", "응답 감사 기록 (운영 누적)"],
               ],
               col_widths=[2.8, 1.8, 7.7], header_size=11, body_size=10)


def slide_data_chroma_neo4j(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "8. 데이터 자산 — Chroma + Neo4j + KIDI + eval",
                    "벡터 DB · 지식 그래프 · 외부 정적 · 평가 셋")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "Chroma 벡터 DB", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(1.0),
               ["컬렉션", "벡터 수", "모델 / 차원"],
               [["insurance_clauses", "739", "text-embedding-3-small / 1536"]],
               col_widths=[2.4, 1.2, 2.4], header_size=11, body_size=10)
    _add_text(s, Inches(0.5), Inches(2.5), Inches(6), Inches(0.4),
              "메타: insurer, product, version, doc_type, clause_no, page",
              size=10, color=COLOR_SUBTLE)

    _add_text(s, Inches(6.8), Inches(1.05), Inches(6), Inches(0.4),
              "Neo4j 지식 그래프 (748 노드)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(2.8),
               ["라벨", "수", "역할"],
               [
                   ["Insurer", "1", "보험사"],
                   ["Product", "2", "상품"],
                   ["Version", "2", "판매기간 버전"],
                   ["Document", "4", "원본 PDF"],
                   ["Clause", "351", "조항 (제N조)"],
                   ["SubClause", "388", "항 (①②③)"],
               ],
               col_widths=[1.6, 0.9, 3.5], header_size=11, body_size=10)
    _add_text(s, Inches(6.8), Inches(4.3), Inches(6), Inches(0.4),
              "엣지 5종: SELLS · HAS_VERSION · HAS_DOCUMENT · CONTAINS · HAS_SUBCLAUSE",
              size=10, color=COLOR_SUBTLE)

    _add_text(s, Inches(0.5), Inches(4.9), Inches(6), Inches(0.4),
              "외부 정적 데이터 + 평가 셋", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6),
               ["데이터셋", "시나리오 수", "내용"],
               [
                   ["KIDI 과실비율 (정적 JSON)", "6", "차101 / 차202 / 차305 / 차411 / 보03 / 이15"],
                   ["eval 시나리오", "10", "eval/scenarios/ — 입력 + 기대 슬롯 + 기대 confidence"],
               ],
               col_widths=[3.5, 1.5, 7.3], header_size=11, body_size=10)


def slide_ops_slo(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "9. 운영 인프라 — SLO + 감사",
                    "Sprint 8: PoC → 대국민 서비스 인프라 전환")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "SLO 목표", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.5),
               ["메트릭", "목표", "측정 위치"],
               [
                   ["API p95 응답시간", "< 5초 (LLM 포함)", "FastAPI 미들웨어"],
                   ["API p50 응답시간", "< 2초", "FastAPI 미들웨어"],
                   ["에러율 (5xx)", "< 0.5% / 24h", "FastAPI 미들웨어"],
                   ["LLM 토큰 비용", "< $0.05 / 응답", "LLM wrapper"],
                   ["RAG 검색 latency p95", "< 1초", "rag_service.retrieve"],
                   ["외부 API 실패율", "< 5% (API별)", "external 어댑터"],
               ],
               col_widths=[3.5, 4.0, 4.8], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(4.3), Inches(6), Inches(0.4),
              "감사 로그 (audit_log)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5),
               ["필드", "내용"],
               [
                   ["response_id", "UUID PK — 분쟁 시 특정 응답 재현 키"],
                   ["session_id", "FK — 세션 삭제 후에도 보존"],
                   ["masked_user_input", "PII 마스킹 후 사용자 입력"],
                   ["llm_calls", "JSONB — 함수명, 모델, 토큰, latency_ms"],
                   ["retrieved_chunk_ids", "text[] — RAG 인용 청크 ID"],
                   ["external_api_calls", "JSONB — 외부 API 호출 기록"],
                   ["보존 기간", "7년 (보험 분쟁 시효 기준, 법무 확인 필수)"],
               ],
               col_widths=[3.0, 9.3], header_size=11, body_size=10)


def slide_ops_pii_breaker(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "9. 운영 인프라 — PII · Rate Limit · 면책",
                    "신뢰성 + 법적 책임 한정")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "PII 마스킹 (입력 · 출력 · 로그 3중)", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4),
               ["대상", "처리"],
               [
                   ["주민번호 · 휴대전화 · 계좌 · 카드 · 이메일", "정규식 + presidio 옵션 마스킹"],
                   ["진단명 · 과실비율 · 사고 경위", "마스킹 제외 (분쟁 시 필수)"],
               ],
               col_widths=[5.0, 7.3], header_size=11, body_size=10)
    _add_text(s, Inches(0.5), Inches(2.95), Inches(12), Inches(0.4),
              "Settings.pii_masking_enabled=False — 테스트 환경 비활성화 가능",
              size=10, color=COLOR_SUBTLE)

    _add_text(s, Inches(0.5), Inches(3.5), Inches(6), Inches(0.4),
              "Rate Limit + Circuit Breaker", size=14, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.6),
               ["위치", "정책"],
               [
                   ["API 진입 (slowapi)", "per-IP 10 req/min · per-session 30 req/min"],
                   ["외부 API 어댑터", "5xx/timeout 5회 연속 → 60초 open → vector RAG 단독 폴백"],
                   ["LLM 호출", "일일 $50 한도 초과 시 503"],
               ],
               col_widths=[3.0, 9.3], header_size=11, body_size=10)

    _add_text(s, Inches(0.5), Inches(5.6), Inches(6), Inches(0.4),
              "면책 + 법적 책임 한정", size=14, bold=True, color=COLOR_WARN)
    _add_table(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.3),
               ["위치", "내용"],
               [
                   ["모든 assessment", "_DEFAULT_DISCLAIMER — '본 결과는 참고용이며 최종 청구 판단을 대체하지 않습니다'"],
                   ["모든 ask + UI 헤더 + /disclaimer 페이지", "면책 영구 표시 + 최초 접속 확인"],
               ],
               col_widths=[4.0, 8.3], header_size=11, body_size=10)


def slide_roadmap(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "10. 로드맵", "Sprint 9 ~ Sprint 12+")

    def block(left, top, width, height, title, status_text, status_color, items):
        _add_box(s, left, top, width, height, fill=COLOR_BG_LIGHT)
        _add_text(s, left + Inches(0.15), top + Inches(0.05), width - Inches(2.0), Inches(0.4),
                  title, size=14, bold=True, color=COLOR_PRIMARY)
        # status badge
        badge = _add_box(s, left + width - Inches(1.9), top + Inches(0.1),
                          Inches(1.7), Inches(0.35),
                          fill=status_color, border=status_color)
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = status_text
        _set_font(r, size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        tb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55),
                                   width - Inches(0.4), height - Inches(0.6))
        tf2 = tb.text_frame; tf2.word_wrap = True
        for i, it in enumerate(items):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.space_after = Pt(2)
            r = p.add_run(); r.text = f"• {it}"
            _set_font(r, size=10, color=COLOR_TEXT)

    block(Inches(0.5), Inches(1.15), Inches(6.1), Inches(2.7),
          "Sprint 9 — 외부 read-only tool 활성",
          "진행 ~60%", COLOR_WARN,
          [
              "KIDI 과실비율 정적 데이터 + get_fault_ratio_standard 활성",
              "lookup_law_clause (법령정보센터) / get_disease_code (HIRA) 어댑터 골격",
              "calc_claim_amount + validate_coverage_period (Sprint 10 선행)",
              "Tool Dispatcher 통합 완료",
              "대기: law OC 코드 · HIRA serviceKey 발급",
          ])

    block(Inches(6.7), Inches(1.15), Inches(6.1), Inches(2.7),
          "Sprint 10 — fss 금감원 크롤링",
          "보류", COLOR_SUBTLE,
          [
              "get_product_meta — 금감원 공시 보험사별 공시실 HTML 스크래핑",
              "각 보험사 공시실 구조 상이 → 복잡도 높음",
              "캐싱 TTL 24시간",
              "Sprint 9 외부 API 완성 후 착수 예정",
          ])

    block(Inches(0.5), Inches(4.0), Inches(6.1), Inches(3.0),
          "Sprint 11 — ReAct agent 본격 활성",
          "진행 ~93%", COLOR_WARN,
          [
              "app/rag/agent.py AgentRunner 구현 (max_iter=5)",
              "rag.service.run_agent 진입점 + sessions.service 분기",
              "audit tool_calls 기록 연동 완료",
              "남은 작업:",
              "  - RAG_REACT=true 환경에서 10 eval 시나리오 전수 통과 확인",
              "  - 영역별 의무/권장 tool 시스템 프롬프트 최종화",
              "  - Grafana 대시보드 연동",
          ])

    block(Inches(6.7), Inches(4.0), Inches(6.1), Inches(3.0),
          "Sprint 12+ — 외부 배포",
          "미착수", COLOR_SUBTLE,
          [
              "도메인 · TLS · DPIA (개인정보영향평가)",
              "클라우드 배포 (운영자 결정)",
              "사용자 인증 (현재 비로그인 → 선택 도입 가능)",
              "Secret Manager 연동 (현재 환경변수)",
              "WCAG AA 접근성 전수 점검 (스크린리더 aria-live 포함)",
          ])


def slide_references(prs):
    s = _slide_blank(prs)
    _add_header_bar(s, "11. 참고 문서 인덱스", "docs/ 하위 분야별 source-of-truth")

    _add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
              "설계 문서 (docs/design/)", size=13, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(0.5), Inches(1.45), Inches(6.2), Inches(5.5),
               ["문서", "내용"],
               [
                   ["tech-decisions.md", "Sprint 1~8.6 기술 결정 + 대안"],
                   ["agent-architecture.md", "LLM agent + ReAct loop + tool 카탈로그"],
                   ["api-spec.md", "CLI + HTTP API + JSON Schema"],
                   ["external-apis.md", "외부 API 4종 (law/HIRA/KIDI/fss)"],
                   ["ui-spec.md", "화면 명세 + 컴포넌트 분해"],
                   ["ui-api-flow.md", "UI ↔ API 흐름 + TS 타입"],
                   ["ui-states.md", "화면 상태 전이"],
                   ["data-model.md", "ERD + SlotState 필드"],
                   ["rag-architecture.md", "Vector/Graph/Hybrid/ReAct 구조"],
                   ["graph-schema.md", "Neo4j 노드/엣지 + Cypher 예시"],
               ],
               col_widths=[2.5, 3.7], header_size=11, body_size=10)

    _add_text(s, Inches(6.9), Inches(1.05), Inches(6), Inches(0.4),
              "사용 가이드 + 요구사항 + PM", size=13, bold=True, color=COLOR_ACCENT)
    _add_table(s, Inches(6.9), Inches(1.45), Inches(6.0), Inches(2.8),
               ["가이드 (docs/)", "내용"],
               [
                   ["usage_sessions.md", "HTTP API curl + ica chat"],
                   ["usage_graphrag.md", "GraphRAG + Neo4j + ica graph-build"],
                   ["usage_response_quality.md", "partial 모드 + 톤 정책"],
                   ["usage_pdf_render.md", "PDF 페이지 캡처 동작"],
                   ["usage_ops.md", "운영자 가이드 — SLO/감사/PII"],
               ],
               col_widths=[2.5, 3.5], header_size=11, body_size=10)

    _add_text(s, Inches(6.9), Inches(4.4), Inches(6), Inches(0.4),
              "요구사항 + 스프린트 + agent", size=13, bold=True, color=COLOR_ACCENT)
    _add_bullet_list(s, Inches(7.1), Inches(4.8), Inches(5.8), Inches(2.4), [
        "docs/requirements/01~09_*.md (REQ-01 ~ REQ-09)",
        "docs/pm/01~11_*.md (스프린트 분석 회의록)",
        "docs/sprint.md — 전체 스프린트 이력",
        "docs/agents/{doc-writer, reviewer, ...}/*.md — agent 보고서",
        "README.md — 빠른 시작 + 환경 변수",
    ], size=11)


def slide_closing(prs):
    s = _slide_blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background(); bg.shadow.inherit = False

    _add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
              "감사합니다",
              size=42, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(0.5),
              "보험청구심사 어시스턴트",
              size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    line.line.fill.background(); line.shadow.inherit = False

    _add_text(s, Inches(1), Inches(4.4), Inches(11.3), Inches(0.5),
              "단일 진입점: docs/SERVICE_OVERVIEW.md",
              size=14, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
              "세부 설계: docs/design/  ·  사용 가이드: docs/usage_*.md",
              size=12, color=RGBColor(0x9C, 0xAA, 0xC0), align=PP_ALIGN.CENTER)

    _add_text(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
              "면책: 본 서비스의 모든 판단 결과는 참고용이며, 최종 보험금 청구 가능 여부의 결정은 보험사에 있습니다.",
              size=10, color=RGBColor(0x9C, 0xAA, 0xC0), align=PP_ALIGN.CENTER)


def build(output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_toc(prs)
    slide_overview(prs)
    slide_progress(prs)
    slide_stack_backend(prs)
    slide_stack_ai(prs)
    slide_stack_data_infra(prs)
    slide_architecture(prs)
    slide_architecture_table(prs)
    slide_api_endpoints(prs)
    slide_api_responses(prs)
    slide_frontend_tree(prs)
    slide_neurosymbolic_neuro(prs)
    slide_neurosymbolic_symbolic(prs)
    slide_scenario_a(prs)
    slide_scenario_b(prs)
    slide_scenario_c(prs)
    slide_scenario_d(prs)
    slide_data_pdf_sql(prs)
    slide_data_chroma_neo4j(prs)
    slide_ops_slo(prs)
    slide_ops_pii_breaker(prs)
    slide_roadmap(prs)
    slide_references(prs)
    slide_closing(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


if __name__ == "__main__":
    target = Path(__file__).resolve().parents[1] / "docs" / "SERVICE_OVERVIEW.pptx"
    saved = build(target)
    print(f"saved: {saved}  ({saved.stat().st_size} bytes)")
