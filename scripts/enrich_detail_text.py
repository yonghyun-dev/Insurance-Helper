"""
슬라이드 3·4의 '세부사항' 텍스트 박스 내용만 풍부하게 교체한다(디자인/박스는 유지).

파일 경로: scripts/enrich_detail_text.py
목적:
    사용자가 디자인을 손본 3·4페이지에서, 우측 세부사항 텍스트만 찾아 내용을 보강한다.
    - AI Agent가 RAG 검색을 어떻게 수행하는지(임베딩·Hybrid·ReRank·동의어·권한필터)
    - 검색 대상(데이터를 어디서 가져오는지: 거버넌스 추출 → Fabrix 임베딩 → Knowledge/Vector DB 인덱싱)
주의:
    텍스트 박스의 글자만 다시 쓰고, 박스 위치·크기·채움/테두리(디자인)는 그대로 둔다.
    넘침 방지를 위해 텍스트 자동 축소(TEXT_TO_FIT_SHAPE)를 켠다.
주요 의존성:
    python-pptx
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

SRC = r"C:\Users\edgar\Desktop\데이터카탈로그 소개서\VODA_AI_Architecture.pptx"

DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x59, 0x59, 0x59)
BLUE = RGBColor(0x2B, 0x5F, 0xC0)
FONT = "맑은 고딕"

# (제목, [하위 bullet]) 블록 리스트 + 끝에 출처 note
DETAIL = {
    # 슬라이드 3: 시나리오 ① 카탈로그 검색
    "검색 질의": {
        "blocks": [
            (
                "검색 질의",
                [
                    "카탈로그 검색 화면에서 자연어·키워드 질의 입력 → AI Agent(Fabrix)에 REST 전달",
                    "사용자 부문·권한 정보 동반 전달",
                ],
            ),
            (
                "AI Agent의 RAG 검색",
                [
                    "질의 임베딩(E5-Large) + 동의어 사전으로 질의 확장(물리명↔논리명·약어)",
                    "Hybrid 검색: Lexical(BM25) + Vector 유사도 동시 수행 → Top-K 후보 추출",
                    "ReRank: 후보(Top-K 10~100)를 재정렬하여 정밀도 향상",
                    "권한·보안등급 메타 필터를 검색 단계에서 적용(부문별 차별 검색)",
                ],
            ),
            (
                "결과 반환",
                [
                    "유사도순 결과를 검색결과 페이지(표·필터)에 표시",
                    "데이터 원본 선택 시 Tableau 리포트 생성 화면으로 이동",
                ],
            ),
        ],
        "note": (
            "검색 대상 (데이터 출처)",
            [
                "거버넌스가 테이블·칼럼 메타(물리·논리명·코멘트·표준분류·보안등급)를 배치 추출",
                "Fabrix 등록 → 청킹·임베딩 → Knowledge/Vector DB에 인덱싱",
                "RAG는 원천 DB가 아닌 이 인덱스(Knowledge)를 조회",
            ],
        ),
    },
    # 슬라이드 4: 시나리오 ② 채팅창 검색
    "자연어 질의": {
        "blocks": [
            (
                "자연어 질의",
                [
                    "AI 서비스 레이어 채팅창에 자연어로 질의 (VODA 포탈 미경유)",
                ],
            ),
            (
                "AI Agent의 질의 해석 · RAG 검색",
                [
                    "Multi-Agent가 질의 의도 해석·분류, 정보 부족 시 질문 구체화 유도",
                    "질의 증강(동의어·물리/논리명) → Hybrid(BM25+Vector) 검색 → ReRank",
                    "권한·보안등급 메타 필터 적용",
                ],
            ),
            (
                "답변 생성",
                [
                    "On-Prem LLM(GPT-OSS 등)이 검색 결과 요약 + 근거(citation) 생성",
                    "채팅창에 대화형 답변 반환",
                ],
            ),
        ],
        "note": (
            "검색 대상 (데이터 출처)",
            [
                "(시나리오 ①과 동일) 거버넌스 추출 메타 → Fabrix 임베딩 → Knowledge/Vector DB 인덱싱",
                "RAG는 이 인덱스를 조회 (원천 DB 직접 접근 아님)",
            ],
        ),
    },
}


def set_run(r, text, size, color, bold):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def rewrite(tf, blocks, note):
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    state = {"first": True}

    def para(text, size, color, bold):
        p = tf.paragraphs[0] if state["first"] else tf.add_paragraph()
        state["first"] = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        r = p.add_run()
        set_run(r, text, size, color, bold)

    for n, (title, subs) in enumerate(blocks, 1):
        para(f"{n}. {title}", 10.5, DARK, True)
        for sb in subs:
            para(f"   • {sb}", 9.5, GRAY, False)
    para("", 4, DARK, False)
    para(f"※ {note[0]}", 10, BLUE, True)
    for sb in note[1]:
        para(f"   • {sb}", 9, GRAY, False)


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)
    # 슬라이드 3=index 2, 슬라이드 4=index 3 만 대상 (아키텍처/다이어그램 슬라이드는 절대 건드리지 않음)
    targets = {2: "검색 질의", 3: "자연어 질의"}
    changed = 0
    for idx, marker in targets.items():
        spec = DETAIL[marker]
        for shape in slides[idx].shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            # 번호 매긴 세부사항 패널만 식별: "1. <marker>" 로 시작
            if txt.startswith(f"1. {marker}"):
                rewrite(shape.text_frame, spec["blocks"], spec["note"])
                changed += 1
                break
    prs.save(SRC)
    print("rewritten boxes:", changed)


if __name__ == "__main__":
    main()
